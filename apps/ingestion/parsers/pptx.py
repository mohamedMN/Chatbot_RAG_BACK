# backend/apps/ingestion/parsers/pptx.py
from typing import List
from io import BytesIO
from pptx import Presentation
from apps.ingestion.schema import Section


def parse_pptx_bytes(data: bytes, filename: str) -> List[Section]:
    prs = Presentation(BytesIO(data))
    out: List[Section] = []
    ordinal = 0
    for i, slide in enumerate(prs.slides, start=1):
        title = None
        if slide.shapes.title and getattr(slide.shapes.title, "text", None):
            title = slide.shapes.title.text.strip()
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    texts.append(t)
        subject = (title or f"Slide {i}")[:120]
        content = "\n".join(texts)
        out.append(Section(
            subject=subject, source=f"{filename}#slide={i}", content=content, ordinal=ordinal))
        ordinal += 1
    return out
